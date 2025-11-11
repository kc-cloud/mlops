#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for Security Controls by Layer
Focuses on the Defense-in-Depth approach with 6 layers
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Security Controls by Layer"
    subtitle.text = "Defense-in-Depth Approach\nSecure MLOps Pipeline for LLM Fine-Tuning"

    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 120, 215)

def add_overview_slide(prs):
    """Add overview slide with all 6 layers"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Defense-in-Depth: 6 Security Layers"

    # Add content
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    layers = [
        ("1", "Model Acquisition Security", "Secure model downloads and integrity verification"),
        ("2", "Container Security", "Vulnerability scanning and hardened containers"),
        ("3", "Training Security", "VPC isolation, encryption, and access controls"),
        ("4", "Data Security", "Encryption at rest, access logging, and retention"),
        ("5", "Model Governance", "Quality gates, versioning, and approval workflows"),
        ("6", "Deployment Security", "Secure endpoints, monitoring, and rate limiting")
    ]

    for num, layer_name, description in layers:
        p = text_frame.add_paragraph()
        p.text = f"Layer {num}: {layer_name}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 120, 215)
        p.space_after = Pt(6)

        p2 = text_frame.add_paragraph()
        p2.text = f"   {description}"
        p2.font.size = Pt(16)
        p2.space_after = Pt(18)

def add_layer_slide(prs, layer_num, layer_name, threat, controls, code_example=None):
    """Add a detailed slide for each security layer"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"Layer {layer_num}: {layer_name}"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 120, 215)

    # Add threat box
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(0.6)
    threat_box = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    threat_box.fill.solid()
    threat_box.fill.fore_color.rgb = RGBColor(220, 50, 50)
    threat_box.line.color.rgb = RGBColor(180, 30, 30)

    text_frame = threat_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = f"⚠ Threat: {threat}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add controls
    left = Inches(0.5)
    top = Inches(2.1)
    width = Inches(9)
    height = Inches(3)

    controls_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = controls_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "Controls Implemented:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 180, 0)
    p.space_after = Pt(12)

    for control in controls:
        p = text_frame.add_paragraph()
        p.text = f"✅ {control}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        p.level = 0

    # Add code example box if provided
    if code_example:
        left = Inches(0.5)
        top = Inches(5.3)
        width = Inches(9)
        height = Inches(1.2)

        code_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = code_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = code_example
        p.font.size = Pt(11)
        p.font.name = 'Courier New'

        # Add background to code box
        code_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(5.3), Inches(9), Inches(1.2)
        )
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = RGBColor(40, 40, 40)
        code_shape.line.color.rgb = RGBColor(80, 80, 80)

        # Move code box to front
        code_shape_elem = code_shape.element
        code_shape_elem.getparent().remove(code_shape_elem)
        slide.shapes._spTree.insert(2, code_shape_elem)

def add_layer1_slide(prs):
    """Layer 1: Model Acquisition Security"""
    controls = [
        "Token-based authentication via AWS Secrets Manager",
        "Model integrity verification (checksums)",
        "Audit logging to CloudWatch",
        "No trust of remote code execution",
        "Secure download over HTTPS only"
    ]

    code = """# Secure token retrieval from Secrets Manager
secrets_client = boto3.client('secretsmanager')
response = secrets_client.get_secret_value(SecretId='huggingface/api-token')
token = response['SecretString']

# Download with verification
model_path = snapshot_download(repo_id=model_id, token=token,
                                local_dir_use_symlinks=False)"""

    add_layer_slide(prs, 1, "Model Acquisition Security",
                   "Malicious model injection, supply chain attacks",
                   controls, code)

def add_layer2_slide(prs):
    """Layer 2: Container Security"""
    controls = [
        "AWS ECR with vulnerability scanning (basic + enhanced)",
        "Immutable image tags",
        "KMS encryption at rest",
        "Non-root container execution",
        "Minimal attack surface (slim base image)",
        "Regular automated scanning",
        "Image signing (optional)"
    ]

    code = """# Dockerfile Hardening
FROM 763104351884.dkr.ecr.us-east-1.amazonaws.com/pytorch-training:2.1.0-gpu-py310
RUN groupadd -g 1000 mlopsuser && useradd -m -u 1000 -g mlopsuser mlopsuser
USER mlopsuser  # Run as non-root"""

    add_layer_slide(prs, 2, "Container Security",
                   "Vulnerable dependencies, container escape, supply chain attacks",
                   controls, code)

def add_layer3_slide(prs):
    """Layer 3: Training Security"""
    controls = [
        "VPC isolation (optional)",
        "Private subnets with no internet access",
        "VPC endpoints for AWS services",
        "Encrypted inter-container traffic",
        "KMS encryption for training data, model artifacts, EBS volumes",
        "IAM least-privilege roles",
        "Network isolation mode"
    ]

    code = """# Security Configuration
security:
  encryption:
    s3_kms_key_id: "alias/sagemaker-kms-key"
    volume_kms_key_id: "alias/sagemaker-volume-kms-key"
    enable_network_isolation: true"""

    add_layer_slide(prs, 3, "Training Security",
                   "Data exfiltration, unauthorized access, model poisoning",
                   controls, code)

def add_layer4_slide(prs):
    """Layer 4: Data Security"""
    controls = [
        "S3 bucket encryption (KMS)",
        "Bucket versioning enabled",
        "Public access blocked",
        "Lifecycle policies for data retention",
        "Access logging",
        "MFA delete protection",
        "Cross-region replication (optional)"
    ]

    code = """# S3 Security Settings
bucket_encryption:
  - SSEAlgorithm: 'aws:kms'
    KMSMasterKeyID: !GetAtt MLOpsKMSKey.Arn
public_access_block:
  BlockPublicAcls: true
  RestrictPublicBuckets: true"""

    add_layer_slide(prs, 4, "Data Security",
                   "Data breach, unauthorized access, compliance violations",
                   controls, code)

def add_layer5_slide(prs):
    """Layer 5: Model Governance"""
    controls = [
        "Automated performance threshold validation",
        "Model versioning in SageMaker Registry",
        "Approval workflow (manual/automated)",
        "Model lineage tracking",
        "Model cards for documentation",
        "Experiment tracking",
        "Audit trail"
    ]

    code = """# Quality Gates
thresholds = {
    'perplexity_max': 20.0,    # Industry best practice
    'eval_loss_max': 1.5,       # Model quality threshold
}
if not check_thresholds(metrics, thresholds):
    approval_status = 'Rejected'"""

    add_layer_slide(prs, 5, "Model Governance",
                   "Unauthorized model deployment, model drift, compliance violations",
                   controls, code)

def add_layer6_slide(prs):
    """Layer 6: Deployment Security"""
    controls = [
        "VPC endpoint deployment",
        "Data capture for monitoring",
        "Auto-scaling with security",
        "HTTPS-only endpoints",
        "IAM-based access control",
        "Request/response encryption",
        "Rate limiting",
        "Model monitoring (drift, bias, quality)"
    ]

    code = """# Endpoint Security
endpoint_config = {
    'DataCaptureConfig': {'EnableCapture': True, 'KmsKeyId': kms_key_id},
    'KmsKeyId': kms_key_id,
    'EnableNetworkIsolation': True
}"""

    add_layer_slide(prs, 6, "Deployment Security",
                   "Unauthorized access, data leakage, model serving attacks",
                   controls, code)

def add_summary_slide(prs):
    """Add summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Security Best Practices Summary"

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    practices = [
        ("Defense in Depth", "Multiple layers of security controls, no single point of failure"),
        ("Least Privilege", "Minimal IAM permissions, service-specific roles"),
        ("Encryption Everywhere", "At rest (S3, EBS, ECR) and in transit (TLS 1.2+)"),
        ("Automated Security", "Automated vulnerability scanning and compliance checks"),
        ("Continuous Monitoring", "Real-time alerts, centralized logging, regular audits")
    ]

    for practice, description in practices:
        p = text_frame.add_paragraph()
        p.text = f"✅ {practice}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 120, 215)
        p.space_after = Pt(6)

        p2 = text_frame.add_paragraph()
        p2.text = f"   {description}"
        p2.font.size = Pt(16)
        p2.space_after = Pt(18)

def add_compliance_slide(prs):
    """Add compliance slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Compliance & Regulatory Standards"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 120, 215)

    # Three columns for different compliance standards
    compliance_data = [
        ("SOC 2 Type II", [
            "Encryption at rest and in transit",
            "Access controls and audit logging",
            "Change management workflows",
            "Monitoring and alerting"
        ]),
        ("HIPAA Ready", [
            "BAA-eligible AWS services",
            "Data encryption",
            "Access logging",
            "Network isolation"
        ]),
        ("GDPR", [
            "Data encryption",
            "Data retention policies",
            "Right to deletion",
            "Complete audit trail"
        ])
    ]

    start_left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(2.8)
    height = Inches(4.5)

    for idx, (standard, requirements) in enumerate(compliance_data):
        left = start_left + (idx * Inches(3.1))

        # Add standard box
        std_box = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, Inches(0.5)
        )
        std_box.fill.solid()
        std_box.fill.fore_color.rgb = RGBColor(0, 120, 215)
        std_box.line.color.rgb = RGBColor(0, 80, 150)

        text_frame = std_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = standard
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Add requirements
        text_box = slide.shapes.add_textbox(left, top + Inches(0.6), width, Inches(3.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for req in requirements:
            p = text_frame.add_paragraph()
            p.text = f"✅ {req}"
            p.font.size = Pt(12)
            p.space_after = Pt(10)

def main():
    """Generate the PowerPoint presentation"""
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Generating presentation...")

    # Add all slides
    add_title_slide(prs)
    print("  ✅ Added title slide")

    add_overview_slide(prs)
    print("  ✅ Added overview slide")

    add_layer1_slide(prs)
    print("  ✅ Added Layer 1: Model Acquisition Security")

    add_layer2_slide(prs)
    print("  ✅ Added Layer 2: Container Security")

    add_layer3_slide(prs)
    print("  ✅ Added Layer 3: Training Security")

    add_layer4_slide(prs)
    print("  ✅ Added Layer 4: Data Security")

    add_layer5_slide(prs)
    print("  ✅ Added Layer 5: Model Governance")

    add_layer6_slide(prs)
    print("  ✅ Added Layer 6: Deployment Security")

    add_summary_slide(prs)
    print("  ✅ Added summary slide")

    add_compliance_slide(prs)
    print("  ✅ Added compliance slide")

    # Save presentation
    output_file = '/Users/mac/git-kc-cloud/mlops/presentation/Security_Controls_by_Layer.pptx'
    prs.save(output_file)

    print(f"\n✅ Presentation saved successfully!")
    print(f"📁 Location: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
